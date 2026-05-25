"""
Extract text from PPTX files and write enriched .txt files to data/scraped/.
Preserves the SOURCE/URL/TOPIC/SECTOR header from the existing .txt file,
then replaces the body with full slide text extracted from the PPTX.

Usage:
    cd backend
    uv run scripts/extract_pptx.py

After running, re-ingest with:
    uv run scripts/ingest_scraped.py --force  (or delete+reingest manually)
"""

import sys
from pathlib import Path
from pptx import Presentation

# Map PPTX filename -> existing scraped .txt filename (for header preservation)
PPTX_TO_TXT = {
    "as-seen-on-tv.pptx": "as-seen-on-tv.txt",
    "signalling-success.pptx": "signalling-success.txt",
    "tv-viewing-report-2025.pptx": "tv-viewing-report-2025.txt",
}


def extract_slide_text(pptx_path: Path, min_chars: int = 20) -> list[str]:
    """
    Extract text from each slide. Returns list of non-empty slide text blocks.
    Skips slides with only whitespace or very short text (titles/footers).
    """
    prs = Presentation(str(pptx_path))
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in para.runs).strip()
                if line:
                    texts.append(line)

        slide_text = "\n".join(texts).strip()
        # Skip slides with negligible content
        if len(slide_text) >= min_chars:
            slides_text.append(f"[Slide {slide_num}]\n{slide_text}")

    return slides_text


def parse_header_lines(txt_path: Path) -> list[str]:
    """Return the header lines (SOURCE/URL/PUBLISHED/TOPIC/SECTOR) from existing .txt."""
    lines = txt_path.read_text(encoding="utf-8").splitlines()
    header = []
    for line in lines:
        if any(line.startswith(k) for k in ("SOURCE:", "URL:", "PUBLISHED:", "TOPIC:", "SECTOR:")):
            header.append(line)
        elif header and line.strip() == "":
            break
    return header


def main():
    backend_dir = Path(__file__).parent.parent
    pptx_dir = backend_dir.parent / "data" / "pptx"
    scraped_dir = backend_dir.parent / "data" / "scraped"

    if not pptx_dir.exists():
        print(f"PPTX directory not found: {pptx_dir}")
        sys.exit(1)

    for pptx_name, txt_name in PPTX_TO_TXT.items():
        pptx_path = pptx_dir / pptx_name
        txt_path = scraped_dir / txt_name

        if not pptx_path.exists():
            print(f"SKIP: {pptx_name} not found")
            continue

        print(f"Extracting: {pptx_name} ...", end=" ", flush=True)
        slides = extract_slide_text(pptx_path)
        print(f"{len(slides)} slides with content")

        # Preserve header from existing .txt if present
        header_lines = parse_header_lines(txt_path) if txt_path.exists() else []

        body = "\n\n".join(slides)
        content = "\n".join(header_lines) + "\n\n" + body if header_lines else body

        txt_path.write_text(content, encoding="utf-8")
        word_count = len(body.split())
        print(f"  Written {word_count} words to {txt_path.name}")


if __name__ == "__main__":
    main()
